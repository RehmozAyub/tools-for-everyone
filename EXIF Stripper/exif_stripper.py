"""
Metadata Scrubber — Inspect & strip metadata from images AND documents
Supports: Images (JPEG, PNG, WebP, TIFF, BMP) + Documents (DOCX, PDF, XLSX, PPTX)
Built with Streamlit + Pillow + piexif + python-docx + openpyxl + python-pptx + PyMuPDF
"""

import streamlit as st
from PIL import Image, ExifTags, ImageOps
from PIL.ExifTags import TAGS, GPSTAGS, IFD
import piexif
import io
import os
import zipfile
import json
import pandas as pd
from datetime import datetime
from pathlib import Path

# Document libraries
from docx import Document as DocxDocument
from docx.opc.constants import RELATIONSHIP_TYPE as RT
import openpyxl
from pptx import Presentation
import fitz  # PyMuPDF

# ─── Page Config ───
st.set_page_config(
    page_title="Metadata Scrubber",
    page_icon="🛡️",
    layout="wide",
    initial_sidebar_state="expanded",
)

# ─── Custom CSS (consistent with PDF Suite / QR Code Studio / Bulk Renamer) ───
st.markdown("""
<style>
    .main-header {
        font-size: 2.5rem;
        font-weight: 700;
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        -webkit-background-clip: text;
        -webkit-text-fill-color: transparent;
        text-align: center;
        margin-bottom: 0.2rem;
    }
    .sub-header {
        text-align: center;
        color: #888;
        margin-bottom: 2rem;
        font-size: 1.05rem;
    }
    .stDownloadButton > button {
        width: 100%;
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        color: white;
        border: none;
        border-radius: 8px;
        padding: 0.6rem 1rem;
        font-weight: 600;
    }
    .stDownloadButton > button:hover {
        background: linear-gradient(135deg, #764ba2 0%, #667eea 100%);
        color: white;
    }
    div[data-testid="stSidebar"] {
        background: #0e1117;
    }
    .tool-info {
        padding: 0.8rem 1rem;
        border-radius: 8px;
        background: #1a1a2e;
        border-left: 4px solid #667eea;
        margin-bottom: 1rem;
    }
    .privacy-warn {
        padding: 0.6rem 1rem;
        border-radius: 8px;
        background: #2a1a1a;
        border-left: 4px solid #e74c3c;
        margin-bottom: 0.5rem;
        color: #e8a0a0;
    }
    .privacy-ok {
        padding: 0.6rem 1rem;
        border-radius: 8px;
        background: #1a2a1a;
        border-left: 4px solid #2ecc71;
        margin-bottom: 0.5rem;
        color: #a0e8a0;
    }
</style>
""", unsafe_allow_html=True)

# ─── Header ───
st.markdown('<div class="main-header">🛡️ Metadata Scrubber</div>', unsafe_allow_html=True)
st.markdown(
    '<div class="sub-header">Inspect & strip metadata from images and documents — protect your privacy</div>',
    unsafe_allow_html=True,
)

# ─── Constants ───
IMAGE_FORMATS = ["jpg", "jpeg", "png", "webp", "tiff", "bmp"]
DOC_FORMATS = ["docx", "pdf", "xlsx", "pptx"]
ALL_FORMATS = IMAGE_FORMATS + DOC_FORMATS

SENSITIVE_IMAGE_TAGS = {
    "GPSLatitude", "GPSLongitude", "GPSLatitudeRef", "GPSLongitudeRef",
    "GPSAltitude", "GPSTimeStamp", "GPSDateStamp",
    "BodySerialNumber", "CameraSerialNumber", "SerialNumber",
    "LensSerialNumber", "ImageUniqueID",
    "CameraOwnerName", "OwnerName", "Artist", "Copyright",
    "XPAuthor", "XPComment",
}

SENSITIVE_DOC_PROPERTIES = {
    "author", "last_modified_by", "creator", "producer",
    "company", "manager", "category",
}


# ═════════════════════════════════════════════════════
# IMAGE HELPERS (existing)
# ═════════════════════════════════════════════════════
def get_exif_tables(img_bytes):
    """Extract all EXIF metadata and return categorized dicts."""
    img = Image.open(io.BytesIO(img_bytes))
    exif = img.getexif()

    basic, camera, gps, advanced, all_tags = {}, {}, {}, {}, {}

    basic["Format"] = img.format or "Unknown"
    basic["Mode"] = img.mode
    basic["Size"] = f"{img.width} × {img.height} px"
    basic["File Size"] = f"{len(img_bytes) / 1024:.1f} KB"

    for tag_id, value in exif.items():
        tag_name = TAGS.get(tag_id, f"Tag_{tag_id}")
        str_value = _safe_str(value)
        all_tags[tag_name] = str_value
        if tag_name in ("Make", "Model", "Software", "LensMake", "LensModel"):
            camera[tag_name] = str_value
        elif tag_name in ("DateTime", "DateTimeOriginal", "DateTimeDigitized"):
            basic[tag_name] = str_value
        elif tag_name == "Orientation":
            basic[tag_name] = _orientation_str(value)
        elif tag_name not in ("ImageWidth", "ImageLength"):
            advanced[tag_name] = str_value

    try:
        exif_ifd = exif.get_ifd(IFD.Exif)
        for tag_id, value in exif_ifd.items():
            tag_name = TAGS.get(tag_id, f"Tag_{tag_id}")
            str_value = _safe_str(value)
            all_tags[tag_name] = str_value
            if tag_name in (
                "ExposureTime", "FNumber", "ISOSpeedRatings",
                "FocalLength", "FocalLengthIn35mmFilm",
                "ExposureProgram", "MeteringMode", "Flash",
                "WhiteBalance", "ExposureBiasValue",
                "MaxApertureValue", "DigitalZoomRatio",
                "ShutterSpeedValue", "ApertureValue",
                "BrightnessValue", "LightSource",
            ):
                camera[tag_name] = str_value
            elif tag_name in ("DateTimeOriginal", "DateTimeDigitized", "OffsetTime", "OffsetTimeOriginal"):
                basic[tag_name] = str_value
            else:
                advanced[tag_name] = str_value
    except Exception:
        pass

    try:
        gps_ifd = exif.get_ifd(IFD.GPSInfo)
        for tag_id, value in gps_ifd.items():
            tag_name = GPSTAGS.get(tag_id, f"GPSTag_{tag_id}")
            str_value = _safe_str(value)
            gps[tag_name] = str_value
            all_tags[tag_name] = str_value
    except Exception:
        pass

    return basic, camera, gps, advanced, all_tags


def parse_gps_coords(img_bytes):
    """Extract GPS coordinates as (lat, lon) or None."""
    try:
        img = Image.open(io.BytesIO(img_bytes))
        exif = img.getexif()
        gps_ifd = exif.get_ifd(IFD.GPSInfo)
        if not gps_ifd:
            return None
        lat = gps_ifd.get(2)
        lat_ref = gps_ifd.get(1)
        lon = gps_ifd.get(4)
        lon_ref = gps_ifd.get(3)
        if not all([lat, lat_ref, lon, lon_ref]):
            return None
        lat_deg = float(lat[0]) + float(lat[1]) / 60.0 + float(lat[2]) / 3600.0
        lon_deg = float(lon[0]) + float(lon[1]) / 60.0 + float(lon[2]) / 3600.0
        if lat_ref == "S":
            lat_deg = -lat_deg
        if lon_ref == "W":
            lon_deg = -lon_deg
        return lat_deg, lon_deg
    except Exception:
        return None


def strip_image_metadata(img_bytes, preserve_orientation=True):
    """Remove all EXIF/metadata from image bytes."""
    img = Image.open(io.BytesIO(img_bytes))
    original_format = img.format or "JPEG"
    if preserve_orientation:
        img = ImageOps.exif_transpose(img) or img

    output_format = original_format.upper()
    if output_format in ("JPG", "JPEG"):
        output_format = "JPEG"
        if img.mode in ("RGBA", "P"):
            img = img.convert("RGB")
    elif output_format not in ("WEBP", "PNG"):
        output_format = "JPEG"
        if img.mode in ("RGBA", "P"):
            img = img.convert("RGB")

    buf = io.BytesIO()
    ext_map = {"JPEG": "jpg", "PNG": "png", "WEBP": "webp"}

    if output_format == "JPEG":
        img.save(buf, format="JPEG", quality=95, optimize=True)
        try:
            return piexif.remove(buf.getvalue()), ext_map.get(output_format, "jpg")
        except Exception:
            return buf.getvalue(), "jpg"
    else:
        img.save(buf, format=output_format, quality=95 if output_format == "WEBP" else None, optimize=True if output_format == "PNG" else False)
        if output_format == "WEBP":
            try:
                return piexif.remove(buf.getvalue()), "webp"
            except Exception:
                pass
        return buf.getvalue(), ext_map.get(output_format, "jpg")


def check_image_privacy_risks(all_tags):
    """Check for privacy-sensitive image tags."""
    return [(t, all_tags[t]) for t in all_tags if t in SENSITIVE_IMAGE_TAGS]


def count_image_metadata_tags(img_bytes):
    try:
        _, _, _, _, all_tags = get_exif_tables(img_bytes)
        return len(all_tags)
    except Exception:
        return 0


def _safe_str(value):
    if isinstance(value, bytes):
        try:
            return value.decode("utf-8", errors="replace").strip("\x00")
        except Exception:
            return f"<{len(value)} bytes>"
    if isinstance(value, tuple) and len(value) == 2 and isinstance(value[1], int):
        if value[1] != 0:
            result = value[0] / value[1]
            return f"{result:.4g}" if result != int(result) else str(int(result))
    return str(value)


def _orientation_str(val):
    return {
        1: "Normal", 2: "Mirrored horizontal", 3: "Rotated 180°",
        4: "Mirrored vertical", 5: "Mirrored horizontal + Rotated 270°",
        6: "Rotated 90° CW", 7: "Mirrored horizontal + Rotated 90°",
        8: "Rotated 270° CW",
    }.get(val, str(val))


# ═════════════════════════════════════════════════════
# DOCUMENT HELPERS (new)
# ═════════════════════════════════════════════════════

def _str_or_none(val):
    """Convert value to string, return None if empty."""
    if val is None:
        return None
    s = str(val).strip()
    return s if s else None


# ---- DOCX ----
def inspect_docx(file_bytes):
    """Extract metadata from a DOCX file. Returns (properties, extras, risks)."""
    doc = DocxDocument(io.BytesIO(file_bytes))
    cp = doc.core_properties

    props = {}
    prop_names = [
        "author", "last_modified_by", "title", "subject", "keywords",
        "category", "comments", "content_status", "version",
        "revision", "created", "modified", "last_printed",
        "identifier", "language",
    ]
    for name in prop_names:
        val = _str_or_none(getattr(cp, name, None))
        if val:
            props[name] = val

    extras = {}
    # Count comments
    try:
        from docx.opc.constants import RELATIONSHIP_TYPE as RT
        comments_part = None
        for rel in doc.part.rels.values():
            if "comments" in rel.reltype:
                comments_part = rel.target_part
                break
        if comments_part:
            from lxml import etree
            tree = etree.fromstring(comments_part.blob)
            comment_count = len(tree)
            if comment_count:
                extras["Comments in document"] = str(comment_count)
    except Exception:
        pass

    # Count sections/paragraphs
    extras["Paragraphs"] = str(len(doc.paragraphs))
    extras["Tables"] = str(len(doc.tables))
    extras["Sections"] = str(len(doc.sections))

    risks = [(k, props[k]) for k in props if k in SENSITIVE_DOC_PROPERTIES]
    return props, extras, risks


def strip_docx(file_bytes):
    """Remove metadata from DOCX. Returns cleaned bytes."""
    doc = DocxDocument(io.BytesIO(file_bytes))
    cp = doc.core_properties

    cp.author = ""
    cp.last_modified_by = ""
    cp.title = ""
    cp.subject = ""
    cp.keywords = ""
    cp.category = ""
    cp.comments = ""
    cp.content_status = ""
    cp.revision = 1

    # Clear app.xml properties (company, manager, etc.)
    try:
        from lxml import etree
        for rel in doc.part.package.rels.values():
            if "extended-properties" in rel.reltype:
                tree = etree.fromstring(rel.target_part.blob)
                ns = tree.nsmap.get(None, "")
                for tag_name in ["Company", "Manager", "Application", "AppVersion", "Template"]:
                    for elem in tree.findall(f"{{{ns}}}{tag_name}"):
                        elem.text = ""
                rel.target_part._blob = etree.tostring(tree, xml_declaration=True, encoding="UTF-8", standalone=True)
                break
    except Exception:
        pass

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


# ---- PDF ----
def inspect_pdf(file_bytes):
    """Extract metadata from PDF. Returns (properties, extras, risks)."""
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    meta = doc.metadata or {}

    props = {}
    for key in ["author", "creator", "producer", "title", "subject", "keywords", "creationDate", "modDate"]:
        val = _str_or_none(meta.get(key))
        if val:
            props[key] = val

    extras = {
        "Pages": str(doc.page_count),
        "Encrypted": str(doc.is_encrypted),
        "Format": meta.get("format", "Unknown"),
    }

    # Check for embedded files
    try:
        embeds = doc.embfile_count()
        if embeds:
            extras["Embedded files"] = str(embeds)
    except Exception:
        pass

    doc.close()

    risks = []
    risk_keys = {"author", "creator", "producer"}
    for k in props:
        if k in risk_keys:
            risks.append((k, props[k]))
    return props, extras, risks


def strip_pdf(file_bytes):
    """Remove metadata from PDF. Returns cleaned bytes."""
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    doc.set_metadata({
        "author": "",
        "creator": "",
        "producer": "",
        "title": "",
        "subject": "",
        "keywords": "",
        "creationDate": "",
        "modDate": "",
    })
    # Scrub XMP metadata
    try:
        doc.del_xml_metadata()
    except Exception:
        pass

    buf = io.BytesIO()
    doc.save(buf, garbage=4, deflate=True, clean=True)
    doc.close()
    return buf.getvalue()


# ---- XLSX ----
def inspect_xlsx(file_bytes):
    """Extract metadata from XLSX. Returns (properties, extras, risks)."""
    wb = openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=False, data_only=True)
    cp = wb.properties

    props = {}
    prop_names = [
        "creator", "lastModifiedBy", "title", "subject", "description",
        "keywords", "category", "version", "created", "modified",
        "last_printed", "company", "manager",
    ]
    for name in prop_names:
        val = _str_or_none(getattr(cp, name, None))
        if val:
            props[name] = val

    extras = {}
    sheet_names = wb.sheetnames
    extras["Sheets"] = str(len(sheet_names))
    extras["Sheet names"] = ", ".join(sheet_names)

    # Detect hidden sheets
    hidden_sheets = []
    for sheet_name in sheet_names:
        ws = wb[sheet_name]
        if ws.sheet_state != "visible":
            hidden_sheets.append(f"{sheet_name} ({ws.sheet_state})")
    if hidden_sheets:
        extras["Hidden sheets"] = ", ".join(hidden_sheets)

    # Count defined names
    if wb.defined_names:
        extras["Named ranges"] = str(len(list(wb.defined_names.definedName)))

    wb.close()

    risks = []
    risk_map = {"creator": "author", "lastModifiedBy": "last_modified_by", "company": "company", "manager": "manager"}
    for k, v in props.items():
        if k in risk_map and risk_map[k] in SENSITIVE_DOC_PROPERTIES:
            risks.append((k, v))
    return props, extras, risks


def strip_xlsx(file_bytes):
    """Remove metadata from XLSX. Returns cleaned bytes."""
    wb = openpyxl.load_workbook(io.BytesIO(file_bytes))
    cp = wb.properties

    cp.creator = ""
    cp.lastModifiedBy = ""
    cp.title = ""
    cp.subject = ""
    cp.description = ""
    cp.keywords = ""
    cp.category = ""
    cp.company = ""
    cp.manager = ""

    buf = io.BytesIO()
    wb.save(buf)
    wb.close()
    return buf.getvalue()


# ---- PPTX ----
def inspect_pptx(file_bytes):
    """Extract metadata from PPTX. Returns (properties, extras, risks)."""
    prs = Presentation(io.BytesIO(file_bytes))
    cp = prs.core_properties

    props = {}
    prop_names = [
        "author", "last_modified_by", "title", "subject", "keywords",
        "category", "comments", "content_status", "version",
        "revision", "created", "modified",
    ]
    for name in prop_names:
        val = _str_or_none(getattr(cp, name, None))
        if val:
            props[name] = val

    extras = {
        "Slides": str(len(prs.slides)),
        "Slide width": str(prs.slide_width),
        "Slide height": str(prs.slide_height),
    }

    # Count speaker notes
    notes_count = 0
    for slide in prs.slides:
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                notes_count += 1
    if notes_count:
        extras["Slides with speaker notes"] = str(notes_count)

    risks = [(k, props[k]) for k in props if k in SENSITIVE_DOC_PROPERTIES]
    return props, extras, risks


def strip_pptx(file_bytes):
    """Remove metadata from PPTX. Returns cleaned bytes."""
    prs = Presentation(io.BytesIO(file_bytes))
    cp = prs.core_properties

    cp.author = ""
    cp.last_modified_by = ""
    cp.title = ""
    cp.subject = ""
    cp.keywords = ""
    cp.category = ""
    cp.comments = ""
    cp.content_status = ""
    cp.revision = 1

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ═════════════════════════════════════════════════════
# UNIFIED HELPERS
# ═════════════════════════════════════════════════════
def get_file_ext(filename):
    return filename.rsplit(".", 1)[-1].lower() if "." in filename else ""


def is_image(filename):
    return get_file_ext(filename) in IMAGE_FORMATS


def is_document(filename):
    return get_file_ext(filename) in DOC_FORMATS


def inspect_document(file_bytes, filename):
    """Inspect any document type. Returns (properties, extras, risks, doc_type)."""
    ext = get_file_ext(filename)
    if ext == "docx":
        props, extras, risks = inspect_docx(file_bytes)
        return props, extras, risks, "Word Document"
    elif ext == "pdf":
        props, extras, risks = inspect_pdf(file_bytes)
        return props, extras, risks, "PDF Document"
    elif ext == "xlsx":
        props, extras, risks = inspect_xlsx(file_bytes)
        return props, extras, risks, "Excel Spreadsheet"
    elif ext == "pptx":
        props, extras, risks = inspect_pptx(file_bytes)
        return props, extras, risks, "PowerPoint Presentation"
    return {}, {}, [], "Unknown"


def strip_document(file_bytes, filename):
    """Strip metadata from any document type. Returns (cleaned_bytes, ext)."""
    ext = get_file_ext(filename)
    if ext == "docx":
        return strip_docx(file_bytes), "docx"
    elif ext == "pdf":
        return strip_pdf(file_bytes), "pdf"
    elif ext == "xlsx":
        return strip_xlsx(file_bytes), "xlsx"
    elif ext == "pptx":
        return strip_pptx(file_bytes), "pptx"
    return file_bytes, ext


def count_doc_metadata(file_bytes, filename):
    try:
        props, _, _, _ = inspect_document(file_bytes, filename)
        return len(props)
    except Exception:
        return 0


FORMAT_ICONS = {
    "docx": "📝", "pdf": "📄", "xlsx": "📊", "pptx": "📽️",
    "jpg": "🖼️", "jpeg": "🖼️", "png": "🖼️", "webp": "🖼️",
    "tiff": "🖼️", "bmp": "🖼️",
}


# ─── Sidebar ───
with st.sidebar:
    st.markdown("### 🧰 Select Tool")
    selected = st.radio(
        "",
        [
            "🔍 Inspect Metadata",
            "🧹 Strip Single File",
            "📦 Batch Strip",
        ],
        label_visibility="collapsed",
    )

    st.markdown("---")

    if selected in ("🧹 Strip Single File", "📦 Batch Strip"):
        st.markdown("### ⚙️ Settings")
        preserve_orient = st.checkbox(
            "Preserve image orientation",
            value=True,
            help="Apply EXIF rotation before stripping so images don’t appear rotated",
        )
    else:
        preserve_orient = True

    st.markdown("---")
    st.markdown("##### 📂 Supported Formats")
    st.markdown(
        "🖼️ Images: JPG, PNG, WebP, TIFF, BMP\n\n"
        "📝 Word: DOCX\n\n"
        "📄 PDF\n\n"
        "📊 Excel: XLSX\n\n"
        "📽️ PowerPoint: PPTX"
    )

    st.markdown("---")
    st.markdown("##### ℹ️ About")
    st.markdown(
        "Free, open-source metadata tool.\n\n"
        "No file uploads to external servers — "
        "**everything runs locally on your machine.**"
    )
    st.markdown("---")
    st.markdown(
        "Made with ❤️ by "
        "[RehmozAyub](https://github.com/RehmozAyub)"
    )


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# INSPECT METADATA
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
if selected == "🔍 Inspect Metadata":
    st.markdown("### 🔍 Inspect Metadata")
    st.markdown("Upload an image or document to view all embedded metadata.")

    file = st.file_uploader(
        "Upload file",
        type=ALL_FORMATS,
        key="inspect_upload",
    )

    if file:
        file_bytes = file.read()
        ext = get_file_ext(file.name)
        icon = FORMAT_ICONS.get(ext, "📁")

        # ---- IMAGE ----
        if is_image(file.name):
            col_img, col_info = st.columns([1, 2])
            with col_img:
                st.image(file_bytes, caption=file.name, width=350)
            with col_info:
                with st.spinner("Reading metadata…"):
                    basic, camera, gps, advanced, all_tags = get_exif_tables(file_bytes)

                risks = check_image_privacy_risks(all_tags)
                if risks:
                    st.markdown(
                        '<div class="privacy-warn">' +
                        f'⚠️ <strong>Privacy Alert:</strong> Found {len(risks)} sensitive tag(s).' +
                        '</div>', unsafe_allow_html=True,
                    )
                    for tag_name, tag_val in risks:
                        st.markdown(f"- ⚠️ **{tag_name}:** `{tag_val}`")
                else:
                    st.markdown(
                        '<div class="privacy-ok">✅ <strong>No sensitive metadata found.</strong></div>',
                        unsafe_allow_html=True,
                    )
                st.info(f"🏷️ **{len(all_tags)}** metadata tags found")

            st.markdown("---")

            coords = parse_gps_coords(file_bytes)
            if coords:
                st.markdown("### 🗺️ GPS Location")
                lat, lon = coords
                st.markdown(f"**Latitude:** `{lat:.6f}` &nbsp;&nbsp; **Longitude:** `{lon:.6f}`")
                st.map(pd.DataFrame({"lat": [lat], "lon": [lon]}), zoom=12)
                st.markdown("---")

            tab_basic, tab_camera, tab_gps, tab_advanced, tab_all = st.tabs(
                ["📋 Basic", "📷 Camera", "📍 GPS", "🔧 Advanced", "📜 All Tags"]
            )
            for tab, data, label in [
                (tab_basic, basic, "basic"), (tab_camera, camera, "camera"),
                (tab_gps, gps, "GPS"), (tab_advanced, advanced, "advanced"),
            ]:
                with tab:
                    if data:
                        st.dataframe(pd.DataFrame(list(data.items()), columns=["Tag", "Value"]), use_container_width=True, hide_index=True)
                    else:
                        st.info(f"No {label} metadata found.")
            with tab_all:
                if all_tags:
                    st.dataframe(pd.DataFrame(list(all_tags.items()), columns=["Tag", "Value"]), use_container_width=True, hide_index=True)
                    st.markdown(f"**Total: {len(all_tags)} tags**")
                else:
                    st.info("No metadata found.")

            if all_tags:
                st.download_button(
                    "⬇️ Export Metadata as JSON",
                    data=json.dumps(all_tags, indent=2, default=str),
                    file_name=f"{file.name}_metadata.json",
                    mime="application/json",
                    use_container_width=True,
                )

        # ---- DOCUMENT ----
        elif is_document(file.name):
            try:
                props, extras, risks, doc_type = inspect_document(file_bytes, file.name)
            except Exception as e:
                st.error(f"❌ Failed to read document: {e}")
                st.stop()

            st.markdown(f"### {icon} {doc_type}")
            st.markdown(f"**File:** {file.name} &nbsp;&nbsp; **Size:** {len(file_bytes)/1024:.1f} KB")

            # Privacy risks
            if risks:
                st.markdown(
                    '<div class="privacy-warn">' +
                    f'⚠️ <strong>Privacy Alert:</strong> Found {len(risks)} identifying property/ies.' +
                    '</div>', unsafe_allow_html=True,
                )
                for prop_name, prop_val in risks:
                    st.markdown(f"- ⚠️ **{prop_name}:** `{prop_val}`")
            else:
                st.markdown(
                    '<div class="privacy-ok">✅ <strong>No sensitive metadata found.</strong></div>',
                    unsafe_allow_html=True,
                )

            st.info(f"🏷️ **{len(props)}** metadata properties found")

            st.markdown("---")

            tab_props, tab_extras = st.tabs(["📋 Properties", "🔧 Document Info"])

            with tab_props:
                if props:
                    st.dataframe(
                        pd.DataFrame(list(props.items()), columns=["Property", "Value"]),
                        use_container_width=True, hide_index=True,
                    )
                else:
                    st.info("No document properties found.")

            with tab_extras:
                if extras:
                    st.dataframe(
                        pd.DataFrame(list(extras.items()), columns=["Info", "Value"]),
                        use_container_width=True, hide_index=True,
                    )
                else:
                    st.info("No extra info.")

            if props:
                st.download_button(
                    "⬇️ Export Metadata as JSON",
                    data=json.dumps({**props, **extras}, indent=2, default=str),
                    file_name=f"{file.name}_metadata.json",
                    mime="application/json",
                    use_container_width=True,
                )


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# STRIP SINGLE FILE
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
elif selected == "🧹 Strip Single File":
    st.markdown("### 🧹 Strip Metadata")
    st.markdown("Remove all metadata from a single image or document.")

    file = st.file_uploader(
        "Upload file",
        type=ALL_FORMATS,
        key="strip_upload",
    )

    if file:
        file_bytes = file.read()
        original_size = len(file_bytes)
        ext = get_file_ext(file.name)
        icon = FORMAT_ICONS.get(ext, "📁")

        col_info1, col_info2 = st.columns([1, 2])

        with col_info1:
            if is_image(file.name):
                st.image(file_bytes, caption=file.name, width=300)
            else:
                st.markdown(f"### {icon}")
                st.markdown(f"**{file.name}**")

        with col_info2:
            st.info(f"📂 Original size: **{original_size / 1024:.1f} KB**")

            if is_image(file.name):
                tag_count = count_image_metadata_tags(file_bytes)
                st.info(f"🏷️ **{tag_count}** metadata tags found")
                coords = parse_gps_coords(file_bytes)
                if coords:
                    st.warning(f"📍 GPS location embedded: `{coords[0]:.4f}, {coords[1]:.4f}`")
                risks = check_image_privacy_risks(get_exif_tables(file_bytes)[4])
                if risks:
                    st.warning(f"⚠️ {len(risks)} privacy-sensitive tag(s) will be removed")
            else:
                try:
                    props, extras, risks, doc_type = inspect_document(file_bytes, file.name)
                    tag_count = len(props)
                    st.info(f"🏷️ **{tag_count}** metadata properties found ({doc_type})")
                    if risks:
                        st.warning(f"⚠️ {len(risks)} identifying properties will be removed")
                        for rk, rv in risks:
                            st.markdown(f"- **{rk}:** `{rv}`")
                except Exception as e:
                    st.error(f"Error reading: {e}")
                    tag_count = 0

        if st.button("🧹 Strip All Metadata", type="primary", use_container_width=True):
            with st.spinner("Stripping metadata…"):
                try:
                    if is_image(file.name):
                        cleaned, out_ext = strip_image_metadata(file_bytes, preserve_orient)
                        new_tags = count_image_metadata_tags(cleaned)
                    else:
                        cleaned, out_ext = strip_document(file_bytes, file.name)
                        new_tags = count_doc_metadata(cleaned, file.name)
                except Exception as e:
                    st.error(f"❌ Stripping failed: {e}")
                    st.stop()

            new_size = len(cleaned)
            reduction = (1 - new_size / original_size) * 100

            c1, c2, c3 = st.columns(3)
            c1.metric("Original", f"{original_size / 1024:.1f} KB")
            c2.metric("Cleaned", f"{new_size / 1024:.1f} KB")
            c3.metric("Properties Remaining", str(new_tags))

            st.success(
                f"✅ Stripped **{max(tag_count - new_tags, 0)}** metadata entries | "
                f"Size change: {reduction:+.1f}%"
            )

            if is_image(file.name):
                col_before, col_after = st.columns(2)
                with col_before:
                    st.markdown("**Before**")
                    st.image(file_bytes, width=300)
                with col_after:
                    st.markdown("**After (cleaned)**")
                    st.image(cleaned, width=300)

            stem = file.name.rsplit(".", 1)[0]
            mime_map = {
                "jpg": "image/jpeg", "png": "image/png", "webp": "image/webp",
                "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                "pdf": "application/pdf",
                "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            }
            st.download_button(
                f"⬇️ Download Cleaned {out_ext.upper()}",
                data=cleaned,
                file_name=f"{stem}_clean.{out_ext}",
                mime=mime_map.get(out_ext, "application/octet-stream"),
                use_container_width=True,
            )


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# BATCH STRIP
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
elif selected == "📦 Batch Strip":
    st.markdown("### 📦 Batch Strip Metadata")
    st.markdown("Upload multiple images and/or documents. Download all cleaned files as ZIP.")

    files = st.file_uploader(
        "Upload files",
        type=ALL_FORMATS,
        accept_multiple_files=True,
        key="batch_upload",
    )

    if files:
        summary = []
        for f in files:
            f_bytes = f.read()
            f.seek(0)
            ext = get_file_ext(f.name)
            icon = FORMAT_ICONS.get(ext, "📁")

            if is_image(f.name):
                tag_count = count_image_metadata_tags(f_bytes)
                has_gps = parse_gps_coords(f_bytes) is not None
                ftype = "Image"
            else:
                tag_count = count_doc_metadata(f_bytes, f.name)
                has_gps = False
                ftype = ext.upper()

            summary.append({
                "Type": f"{icon} {ftype}",
                "File": f.name,
                "Size (KB)": f"{f.size / 1024:.1f}",
                "Metadata": tag_count,
                "GPS": "📍" if has_gps else "—",
            })

        df_summary = pd.DataFrame(summary)
        st.dataframe(df_summary, use_container_width=True, hide_index=True)

        total_meta = sum(s["Metadata"] for s in summary)
        img_count = sum(1 for f in files if is_image(f.name))
        doc_count = sum(1 for f in files if is_document(f.name))

        st.info(
            f"📦 **{len(files)}** files ({img_count} images, {doc_count} documents) | "
            f"🏷️ **{total_meta}** total metadata entries"
        )

        if st.button("🧹 Strip All Files", type="primary", use_container_width=True):
            with st.spinner(f"Stripping metadata from {len(files)} files…"):
                zbuf = io.BytesIO()
                total_stripped = 0
                total_original = 0
                total_cleaned_size = 0

                with zipfile.ZipFile(zbuf, "w", zipfile.ZIP_DEFLATED) as zf:
                    for f in files:
                        raw = f.read()
                        total_original += len(raw)

                        try:
                            if is_image(f.name):
                                old_count = count_image_metadata_tags(raw)
                                cleaned, out_ext = strip_image_metadata(raw, preserve_orient)
                                new_count = count_image_metadata_tags(cleaned)
                            else:
                                old_count = count_doc_metadata(raw, f.name)
                                cleaned, out_ext = strip_document(raw, f.name)
                                new_count = count_doc_metadata(cleaned, f.name)

                            total_stripped += max(old_count - new_count, 0)
                            total_cleaned_size += len(cleaned)

                            stem = f.name.rsplit(".", 1)[0]
                            zf.writestr(f"{stem}_clean.{out_ext}", cleaned)
                        except Exception as e:
                            st.warning(f"⚠️ Skipped {f.name}: {e}")
                            zf.writestr(f.name, raw)  # Include original
                            total_cleaned_size += len(raw)

            reduction = (1 - total_cleaned_size / total_original) * 100 if total_original else 0

            c1, c2, c3 = st.columns(3)
            c1.metric("Files Processed", str(len(files)))
            c2.metric("Metadata Stripped", str(total_stripped))
            c3.metric("Size Change", f"{reduction:+.1f}%")

            st.success(f"✅ Cleaned {len(files)} files | Removed {total_stripped} metadata entries")

            st.download_button(
                "⬇️ Download All (ZIP)",
                data=zbuf.getvalue(),
                file_name=f"metadata_stripped_{datetime.now().strftime('%Y%m%d_%H%M%S')}.zip",
                mime="application/zip",
                use_container_width=True,
            )


# ─── Footer ───
st.markdown("---")
st.markdown(
    '<div style="text-align:center; color:#666; font-size:0.85rem;">'
    'Metadata Scrubber v2.0 · Built with Streamlit + Pillow + piexif + python-docx + openpyxl + python-pptx + PyMuPDF · '
    '<a href="https://github.com/RehmozAyub/tools-for-everyone" style="color:#667eea;">GitHub</a>'
    "</div>",
    unsafe_allow_html=True,
)
